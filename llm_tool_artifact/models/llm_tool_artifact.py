# -*- coding: utf-8 -*-
import base64
import io
import json
import logging
import re
from typing import Any, Optional

from odoo import _, api, models
from odoo.exceptions import UserError

_logger = logging.getLogger(__name__)


class LLMToolArtifact(models.Model):
    _inherit = "llm.tool"

    @api.model
    def _get_available_implementations(self):
        impl = super()._get_available_implementations()
        return impl + [(
            "llm_artifact_builder",
            "Generador de artefactos (gráfico, Excel, PDF, Word, PowerPoint)",
        )]

    # ==================================================================
    # Entry point
    # ==================================================================
    def llm_artifact_builder_execute(
        self,
        artifact_type: str = "echarts",
        title: str = "Informe",
        # --- ECharts (modo principal) ---
        echarts_option: str = "{}",
        explanation: str = "",
        # --- Excel / PDF / Word / PowerPoint ---
        file_name: Optional[str] = None,
        data_json: str = "[]",
        # --- Documentos enriquecidos (PDF / Word / PowerPoint) ---
        document_json: str = "{}",
        # --- Matplotlib legacy (fallback) ---
        chart_kind: str = "bar",
    ) -> dict[str, Any]:
        """Genera artefactos visuales y documentos descargables en el chat.

        ``artifact_type`` admite los siguientes valores:

        * ``"echarts"`` (por defecto) — gráfico interactivo con Apache ECharts.
        * ``"xlsx"``   — Excel descargable. Acepta tanto una lista plana
          como un dict con varias hojas (``sheets``), encabezados, totales,
          columnas calculadas y formato condicional.
        * ``"pdf"``    — Documento PDF rico (portada, secciones, tablas,
          listas, imágenes y citas) generado con ReportLab.
        * ``"docx"``   — Documento Word con estilos (encabezados, párrafos,
          listas, tablas, imágenes) generado con python-docx.
        * ``"pptx"``   — Presentación PowerPoint con varias diapositivas
          (portada, contenido, listas, tablas e imágenes) generada con
          python-pptx.
        * ``"chart"``  — PNG con matplotlib (modo legacy / fallback).

        IMPORTANTE: para los gráficos ECharts, el campo ``mensaje_markdown``
        del resultado debe copiarse LITERALMENTE en tu respuesta para que el
        gráfico aparezca en el chat. Para xlsx/pdf/docx/pptx, comparte el
        ``url_descarga`` para que el usuario pueda bajar el archivo.
        """
        self.ensure_one()
        atype = (artifact_type or "echarts").lower().strip()

        if atype == "xlsx":
            data = self._safe_loads(data_json, default=[], param="data_json")
            doc = self._safe_loads(document_json, default={}, param="document_json")
            return self._artifact_xlsx(
                title=title,
                data=data,
                document=doc,
                fname=file_name or "artefacto",
            )

        if atype == "pdf":
            doc = self._safe_loads(document_json, default={}, param="document_json")
            return self._artifact_pdf(
                title=title, document=doc, fname=file_name or "documento"
            )

        if atype == "docx":
            doc = self._safe_loads(document_json, default={}, param="document_json")
            return self._artifact_docx(
                title=title, document=doc, fname=file_name or "documento"
            )

        if atype == "pptx":
            doc = self._safe_loads(document_json, default={}, param="document_json")
            return self._artifact_pptx(
                title=title, document=doc, fname=file_name or "presentacion"
            )

        if atype == "chart":
            data = self._safe_loads(data_json, default=[], param="data_json")
            return self._artifact_chart(
                title=title,
                data=data,
                chart_kind=chart_kind,
                fname=file_name or "grafico",
            )

        # Default: ECharts interactivo
        return self._artifact_echarts(title, echarts_option, explanation)

    # ==================================================================
    # Helpers
    # ==================================================================
    @staticmethod
    def _safe_loads(raw, default=None, param="json"):
        if not raw or (isinstance(raw, str) and raw.strip() in ("", "{}", "[]")):
            return default if default is not None else {}
        if isinstance(raw, (dict, list)):
            return raw
        try:
            return json.loads(raw)
        except json.JSONDecodeError as e:
            raise UserError(
                _("JSON inválido en %(param)s: %(err)s") % {"param": param, "err": e}
            ) from e

    @staticmethod
    def _slugify(name: str, default: str = "documento") -> str:
        cleaned = re.sub(r"[^A-Za-z0-9_\-]+", "_", (name or "").strip())
        cleaned = cleaned.strip("_")
        return cleaned or default

    def _store_attachment(self, fname: str, content: bytes, mimetype: str) -> dict:
        att = self.env["ir.attachment"].create(
            {
                "name": fname,
                "type": "binary",
                "datas": base64.b64encode(content).decode(),
                "mimetype": mimetype,
            }
        )
        base = self.env["ir.config_parameter"].sudo().get_param("web.base.url", "")
        return {
            "url_descarga": f"{base}/web/content/{att.id}?download=true",
            "attachment_id": att.id,
            "nombre_archivo": fname,
        }

    # ==================================================================
    # ECharts (motor principal)
    # ==================================================================
    def _artifact_echarts(self, title, echarts_option_str, explanation):
        if not echarts_option_str or str(echarts_option_str).strip() in ("{}", ""):
            raise UserError(
                _(
                    "echarts_option está vacío. Proporciona un JSON de "
                    "opción ECharts válido con title, xAxis/yAxis (si "
                    "aplica) y series."
                )
            )
        try:
            option = json.loads(echarts_option_str)
        except json.JSONDecodeError as e:
            raise UserError(_("JSON inválido en echarts_option: %s") % e) from e

        if not isinstance(option, dict):
            raise UserError(_("echarts_option debe ser un objeto JSON (dict)."))

        if "title" not in option and title:
            option["title"] = {"text": title}

        option.setdefault("tooltip", {"trigger": "axis"})
        option.setdefault(
            "toolbox",
            {
                "show": True,
                "feature": {
                    "dataZoom": {"yAxisIndex": "none"},
                    "restore": {},
                    "saveAsImage": {},
                },
            },
        )

        option_str = json.dumps(option, ensure_ascii=False)
        explanation_md = (
            f"\n\n{explanation}" if explanation and explanation.strip() else ""
        )
        mensaje_markdown = f"```echarts\n{option_str}\n```{explanation_md}"

        return {
            "tipo": "echarts",
            "mensaje_markdown": mensaje_markdown,
            "nota": (
                "Copia el campo 'mensaje_markdown' literalmente en tu respuesta "
                "para que el gráfico interactivo aparezca en el chat."
            ),
        }

    # ==================================================================
    # Excel multi-hoja
    # ==================================================================
    def _artifact_xlsx(self, title, data, document, fname):
        """Excel con soporte de varias hojas, encabezados, totales y formato.

        Acepta dos formatos en ``document_json``:

        * ``{"sheets": [ {"name": "Resumen", "rows": [...]}, ... ]}``
        * Vacío y ``data_json`` con la forma clásica
          (lista de dicts → una sola hoja).

        Cada hoja puede definir:

        * ``name``: nombre de la pestaña (máx. 31 caracteres).
        * ``headers``: lista ordenada de columnas. Si no se da, se infiere
          de las claves del primer registro.
        * ``rows``: lista de dicts con los datos (o lista de listas).
        * ``column_widths``: dict ``{"campo": ancho}`` opcional.
        * ``totals``: lista de columnas a totalizar al final con SUM.
        * ``freeze_header``: bool (default True).
        * ``autofilter``: bool (default True).
        * ``zebra``: bool — alternar color por fila (default True).
        """
        try:
            import xlsxwriter
        except ImportError as e:
            raise UserError(_("Instale xlsxwriter: %s") % e) from e

        # Normalizar a estructura multi-hoja
        sheets = []
        if isinstance(document, dict) and document.get("sheets"):
            sheets = document["sheets"]
        elif isinstance(data, list) and data:
            sheets = [{"name": title or "Hoja1", "rows": data}]
        elif isinstance(data, dict) and data.get("sheets"):
            sheets = data["sheets"]
        else:
            sheets = [{"name": title or "Hoja1", "rows": []}]

        buf = io.BytesIO()
        workbook = xlsxwriter.Workbook(buf, {"in_memory": True})

        # Formatos comunes
        fmt_header = workbook.add_format({
            "bold": True,
            "bg_color": "#305496",
            "font_color": "white",
            "border": 1,
            "align": "center",
            "valign": "vcenter",
        })
        fmt_row = workbook.add_format({"border": 1})
        fmt_zebra = workbook.add_format({"border": 1, "bg_color": "#F2F2F2"})
        fmt_total = workbook.add_format({
            "bold": True,
            "top": 2,
            "bg_color": "#FCE4D6",
            "border": 1,
        })
        fmt_money = workbook.add_format(
            {"border": 1, "num_format": '#,##0.00 "USD"'}
        )

        used_names = set()
        for idx, sheet_def in enumerate(sheets):
            raw_name = (
                sheet_def.get("name") if isinstance(sheet_def, dict) else None
            ) or f"Hoja{idx + 1}"
            name = re.sub(r"[\[\]:*?/\\]", "_", str(raw_name))[:31] or f"Hoja{idx + 1}"
            base = name
            n = 1
            while name in used_names:
                n += 1
                suffix = f"_{n}"
                name = (base[: 31 - len(suffix)] + suffix)
            used_names.add(name)
            sheet = workbook.add_worksheet(name)

            rows = sheet_def.get("rows", []) if isinstance(sheet_def, dict) else []
            if not isinstance(rows, list):
                rows = []

            headers = (
                sheet_def.get("headers")
                if isinstance(sheet_def, dict)
                else None
            )
            if not headers:
                if rows and isinstance(rows[0], dict):
                    headers = list(rows[0].keys())
                else:
                    headers = []

            zebra = sheet_def.get("zebra", True) if isinstance(sheet_def, dict) else True
            autofilter = (
                sheet_def.get("autofilter", True)
                if isinstance(sheet_def, dict)
                else True
            )
            freeze_header = (
                sheet_def.get("freeze_header", True)
                if isinstance(sheet_def, dict)
                else True
            )
            money_cols = set(
                sheet_def.get("money_columns", [])
                if isinstance(sheet_def, dict)
                else []
            )
            column_widths = (
                sheet_def.get("column_widths", {})
                if isinstance(sheet_def, dict)
                else {}
            )
            totals = (
                sheet_def.get("totals", [])
                if isinstance(sheet_def, dict)
                else []
            )
            description = (
                sheet_def.get("description", "")
                if isinstance(sheet_def, dict)
                else ""
            )

            row_offset = 0
            if description:
                sheet.merge_range(
                    0,
                    0,
                    0,
                    max(len(headers) - 1, 0),
                    description,
                    workbook.add_format(
                        {"italic": True, "font_color": "#404040", "align": "left"}
                    ),
                )
                row_offset = 1

            # Encabezados
            for c, h in enumerate(headers):
                sheet.write(row_offset, c, str(h), fmt_header)
                width = column_widths.get(h)
                if isinstance(width, (int, float)) and width > 0:
                    sheet.set_column(c, c, width)
                else:
                    sheet.set_column(c, c, max(12, min(40, len(str(h)) + 4)))

            # Filas
            for r, row in enumerate(rows, start=row_offset + 1):
                row_format = fmt_zebra if (zebra and (r - row_offset) % 2 == 0) else fmt_row
                if isinstance(row, dict):
                    for c, h in enumerate(headers):
                        value = row.get(h)
                        cell_fmt = fmt_money if h in money_cols else row_format
                        sheet.write(r, c, value, cell_fmt)
                elif isinstance(row, list):
                    for c, value in enumerate(row):
                        sheet.write(r, c, value, row_format)
                else:
                    sheet.write(r, 0, str(row), row_format)

            # Totales
            if totals and headers:
                total_row = row_offset + 1 + len(rows)
                sheet.write(total_row, 0, "TOTAL", fmt_total)
                for c, h in enumerate(headers):
                    if c == 0:
                        continue
                    if h in totals and len(rows):
                        first = row_offset + 2  # 1-indexed in Excel
                        last = row_offset + 1 + len(rows)
                        col_letter = self._excel_col_letter(c)
                        sheet.write_formula(
                            total_row,
                            c,
                            f"=SUM({col_letter}{first}:{col_letter}{last})",
                            fmt_total,
                        )
                    else:
                        sheet.write(total_row, c, "", fmt_total)

            # Filtros + congelar fila
            if headers:
                sheet.autofilter(
                    row_offset, 0, row_offset + len(rows), len(headers) - 1
                ) if autofilter else None
                if freeze_header:
                    sheet.freeze_panes(row_offset + 1, 0)

            if not headers and not rows:
                sheet.write(0, 0, "Hoja sin datos", fmt_header)

        workbook.close()
        slug = self._slugify(fname, "artefacto")
        return {
            "tipo": "xlsx",
            **self._store_attachment(
                f"{slug}.xlsx",
                buf.getvalue(),
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            ),
            "nota": "Excel multi-hoja generado. Comparte 'url_descarga' con el usuario.",
        }

    @staticmethod
    def _excel_col_letter(col_index_zero_based: int) -> str:
        """Convierte 0→A, 1→B, …, 25→Z, 26→AA…"""
        n = col_index_zero_based + 1
        s = ""
        while n > 0:
            n, r = divmod(n - 1, 26)
            s = chr(65 + r) + s
        return s

    # ==================================================================
    # PDF (ReportLab)
    # ==================================================================
    def _artifact_pdf(self, title, document, fname):
        """Genera un PDF rico con secciones, tablas, listas, imágenes…

        Estructura esperada de ``document_json`` (todo opcional salvo
        ``sections``):

        ```json
        {
          "subtitle": "Reporte mensual",
          "author": "Asistente IA",
          "cover": true,
          "toc": false,
          "sections": [
            {"type": "heading", "level": 1, "text": "Introducción"},
            {"type": "paragraph", "text": "Texto largo…"},
            {"type": "list", "items": ["Punto 1", "Punto 2"], "ordered": false},
            {"type": "table", "headers": ["A","B"], "rows": [[1,2],[3,4]],
             "title": "Resumen"},
            {"type": "image", "data_b64": "...", "caption": "Gráfico"},
            {"type": "quote", "text": "Cita textual."},
            {"type": "spacer", "height": 12},
            {"type": "page_break"}
          ]
        }
        ```
        """
        try:
            from reportlab.lib import colors
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import cm
            from reportlab.platypus import (
                Image,
                ListFlowable,
                ListItem,
                PageBreak,
                Paragraph,
                SimpleDocTemplate,
                Spacer,
                Table,
                TableStyle,
            )
        except ImportError as e:
            raise UserError(
                _(
                    "Instale reportlab para generar PDFs: %s "
                    "(`pip install reportlab`)."
                )
                % e
            ) from e

        sections = []
        if isinstance(document, dict):
            sections = document.get("sections", []) or []
        if not sections:
            raise UserError(
                _(
                    "document_json debe incluir 'sections' (lista de bloques) "
                    "para generar un PDF."
                )
            )

        buf = io.BytesIO()
        doc = SimpleDocTemplate(
            buf,
            pagesize=A4,
            title=title or "Documento",
            author=document.get("author") or "Odoo AI",
            leftMargin=2 * cm,
            rightMargin=2 * cm,
            topMargin=2 * cm,
            bottomMargin=2 * cm,
        )
        styles = getSampleStyleSheet()
        styles.add(
            ParagraphStyle(
                name="OdooQuote",
                parent=styles["Italic"],
                leftIndent=20,
                textColor=colors.HexColor("#555555"),
                spaceBefore=6,
                spaceAfter=6,
                borderPadding=6,
            )
        )

        story = []
        if document.get("cover", True):
            story.append(Spacer(1, 4 * cm))
            story.append(
                Paragraph(
                    f"<para alignment='center'><b>{title or 'Documento'}</b></para>",
                    ParagraphStyle(
                        name="OdooCover",
                        parent=styles["Title"],
                        fontSize=28,
                        leading=34,
                    ),
                )
            )
            if document.get("subtitle"):
                story.append(Spacer(1, 0.5 * cm))
                story.append(
                    Paragraph(
                        f"<para alignment='center'>{document['subtitle']}</para>",
                        styles["Heading2"],
                    )
                )
            if document.get("author"):
                story.append(Spacer(1, 0.5 * cm))
                story.append(
                    Paragraph(
                        f"<para alignment='center'><i>{document['author']}</i></para>",
                        styles["Normal"],
                    )
                )
            story.append(PageBreak())

        for block in sections:
            if not isinstance(block, dict):
                continue
            btype = (block.get("type") or "paragraph").lower()
            if btype == "heading":
                level = max(1, min(int(block.get("level", 1) or 1), 4))
                story.append(
                    Paragraph(
                        str(block.get("text", "")),
                        styles[f"Heading{level}"],
                    )
                )
            elif btype == "paragraph":
                story.append(Paragraph(str(block.get("text", "")), styles["BodyText"]))
                story.append(Spacer(1, 0.2 * cm))
            elif btype == "quote":
                story.append(Paragraph(str(block.get("text", "")), styles["OdooQuote"]))
            elif btype == "list":
                items = block.get("items") or []
                bullet = "1" if block.get("ordered") else "bullet"
                lf = ListFlowable(
                    [
                        ListItem(Paragraph(str(it), styles["BodyText"]))
                        for it in items
                    ],
                    bulletType=bullet,
                )
                story.append(lf)
                story.append(Spacer(1, 0.2 * cm))
            elif btype == "table":
                headers = block.get("headers") or []
                rows = block.get("rows") or []
                table_data = []
                if headers:
                    table_data.append(list(headers))
                for r in rows:
                    if isinstance(r, dict) and headers:
                        table_data.append([r.get(h, "") for h in headers])
                    elif isinstance(r, list):
                        table_data.append(r)
                if table_data:
                    if block.get("title"):
                        story.append(
                            Paragraph(
                                f"<b>{block['title']}</b>", styles["Heading4"]
                            )
                        )
                    t = Table(table_data, repeatRows=1 if headers else 0)
                    style_cmds = [
                        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#888")),
                        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                        ("FONTSIZE", (0, 0), (-1, -1), 9),
                    ]
                    if headers:
                        style_cmds += [
                            (
                                "BACKGROUND",
                                (0, 0),
                                (-1, 0),
                                colors.HexColor("#305496"),
                            ),
                            ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
                            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                        ]
                        # Zebra rows
                        for ridx in range(1, len(table_data)):
                            if ridx % 2 == 0:
                                style_cmds.append(
                                    (
                                        "BACKGROUND",
                                        (0, ridx),
                                        (-1, ridx),
                                        colors.HexColor("#F2F2F2"),
                                    )
                                )
                    t.setStyle(TableStyle(style_cmds))
                    story.append(t)
                    story.append(Spacer(1, 0.3 * cm))
            elif btype == "image":
                b64 = block.get("data_b64") or block.get("image_b64")
                if b64:
                    try:
                        img_bytes = base64.b64decode(b64)
                        img_buf = io.BytesIO(img_bytes)
                        max_w = doc.width
                        story.append(Image(img_buf, width=max_w, height=max_w * 0.6))
                        if block.get("caption"):
                            story.append(
                                Paragraph(
                                    f"<i>{block['caption']}</i>",
                                    styles["Italic"],
                                )
                            )
                        story.append(Spacer(1, 0.3 * cm))
                    except Exception as err:
                        _logger.warning("PDF image inválida: %s", err)
            elif btype == "spacer":
                h = float(block.get("height", 12) or 12)
                story.append(Spacer(1, h))
            elif btype == "page_break":
                story.append(PageBreak())
            else:
                story.append(Paragraph(str(block.get("text", "")), styles["BodyText"]))

        doc.build(story)
        slug = self._slugify(fname, "documento")
        return {
            "tipo": "pdf",
            **self._store_attachment(
                f"{slug}.pdf",
                buf.getvalue(),
                "application/pdf",
            ),
            "nota": "PDF generado. Comparte 'url_descarga' con el usuario.",
        }

    # ==================================================================
    # Word (python-docx)
    # ==================================================================
    def _artifact_docx(self, title, document, fname):
        """Genera un .docx con encabezados, párrafos, tablas, listas e imágenes.

        Misma estructura ``sections`` que el PDF.
        """
        try:
            from docx import Document
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.shared import Cm, Pt, RGBColor
        except ImportError as e:
            raise UserError(
                _(
                    "Instale python-docx para generar Word: %s "
                    "(`pip install python-docx`)."
                )
                % e
            ) from e

        sections = []
        if isinstance(document, dict):
            sections = document.get("sections", []) or []
        if not sections:
            raise UserError(
                _(
                    "document_json debe incluir 'sections' (lista de bloques) "
                    "para generar un Word."
                )
            )

        d = Document()
        # Portada
        if document.get("cover", True):
            t = d.add_paragraph()
            t.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = t.add_run(title or "Documento")
            run.bold = True
            run.font.size = Pt(28)
            run.font.color.rgb = RGBColor(0x30, 0x54, 0x96)
            if document.get("subtitle"):
                p = d.add_paragraph()
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                p.add_run(document["subtitle"]).italic = True
            if document.get("author"):
                p = d.add_paragraph()
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                p.add_run(document["author"])
            d.add_page_break()

        for block in sections:
            if not isinstance(block, dict):
                continue
            btype = (block.get("type") or "paragraph").lower()
            if btype == "heading":
                level = max(1, min(int(block.get("level", 1) or 1), 4))
                d.add_heading(str(block.get("text", "")), level=level)
            elif btype == "paragraph":
                d.add_paragraph(str(block.get("text", "")))
            elif btype == "quote":
                p = d.add_paragraph()
                run = p.add_run(str(block.get("text", "")))
                run.italic = True
                run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
            elif btype == "list":
                items = block.get("items") or []
                style = "List Number" if block.get("ordered") else "List Bullet"
                for it in items:
                    d.add_paragraph(str(it), style=style)
            elif btype == "table":
                headers = block.get("headers") or []
                rows = block.get("rows") or []
                if not headers and rows and isinstance(rows[0], dict):
                    headers = list(rows[0].keys())
                ncols = len(headers) if headers else (
                    len(rows[0]) if rows and isinstance(rows[0], list) else 1
                )
                table = d.add_table(rows=1 if headers else 0, cols=max(1, ncols))
                table.style = "Light Grid Accent 1"
                if headers:
                    hdr = table.rows[0].cells
                    for i, h in enumerate(headers):
                        hdr[i].text = str(h)
                        for run in hdr[i].paragraphs[0].runs:
                            run.bold = True
                for r in rows:
                    cells = table.add_row().cells
                    if isinstance(r, dict) and headers:
                        for i, h in enumerate(headers):
                            cells[i].text = str(r.get(h, ""))
                    elif isinstance(r, list):
                        for i, v in enumerate(r):
                            if i < len(cells):
                                cells[i].text = str(v)
                if block.get("title"):
                    p = d.add_paragraph()
                    p.add_run(block["title"]).bold = True
            elif btype == "image":
                b64 = block.get("data_b64") or block.get("image_b64")
                if b64:
                    try:
                        img_bytes = base64.b64decode(b64)
                        img_buf = io.BytesIO(img_bytes)
                        d.add_picture(img_buf, width=Cm(15))
                        if block.get("caption"):
                            p = d.add_paragraph()
                            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                            p.add_run(block["caption"]).italic = True
                    except Exception as err:
                        _logger.warning("DOCX image inválida: %s", err)
            elif btype == "spacer":
                d.add_paragraph("")
            elif btype == "page_break":
                d.add_page_break()
            else:
                d.add_paragraph(str(block.get("text", "")))

        buf = io.BytesIO()
        d.save(buf)
        slug = self._slugify(fname, "documento")
        return {
            "tipo": "docx",
            **self._store_attachment(
                f"{slug}.docx",
                buf.getvalue(),
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ),
            "nota": "Word generado. Comparte 'url_descarga' con el usuario.",
        }

    # ==================================================================
    # PowerPoint (python-pptx)
    # ==================================================================
    def _artifact_pptx(self, title, document, fname):
        """Genera un .pptx con varias diapositivas.

        Estructura esperada de ``document_json``:

        ```json
        {
          "subtitle": "Resumen ejecutivo Q1",
          "slides": [
            {"layout": "title",   "title": "Resumen Q1", "subtitle": "Equipo Ventas"},
            {"layout": "content", "title": "Logros",     "bullets": ["..."]},
            {"layout": "table",   "title": "Por producto",
             "headers": ["Producto","Ventas"], "rows": [["A",100]]},
            {"layout": "image",   "title": "Gráfico",    "data_b64": "..."}
          ]
        }
        ```
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
        except ImportError as e:
            raise UserError(
                _(
                    "Instale python-pptx para generar PowerPoint: %s "
                    "(`pip install python-pptx`)."
                )
                % e
            ) from e

        slides = []
        if isinstance(document, dict):
            slides = document.get("slides", []) or []
        if not slides:
            # Si no hay slides, creamos una portada sencilla.
            slides = [
                {
                    "layout": "title",
                    "title": title or "Presentación",
                    "subtitle": (document or {}).get("subtitle") or "",
                }
            ]

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for slide_def in slides:
            if not isinstance(slide_def, dict):
                continue
            layout = (slide_def.get("layout") or "content").lower()
            slide_layout = prs.slide_layouts[5]  # Blank
            slide = prs.slides.add_slide(slide_layout)

            # Título superior
            if slide_def.get("title"):
                tx = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9)
                )
                tf = tx.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = str(slide_def["title"])
                run.font.size = Pt(32)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x30, 0x54, 0x96)

            if layout == "title":
                tx = slide.shapes.add_textbox(
                    Inches(0.5), Inches(3.0), Inches(12.3), Inches(2.0)
                )
                tf = tx.text_frame
                tf.word_wrap = True
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                run = tf.paragraphs[0].add_run()
                run.text = str(slide_def.get("subtitle") or "")
                run.font.size = Pt(20)
                continue

            if layout == "content" or slide_def.get("bullets"):
                bullets = slide_def.get("bullets") or [
                    slide_def.get("text", "")
                ]
                tx = slide.shapes.add_textbox(
                    Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5)
                )
                tf = tx.text_frame
                tf.word_wrap = True
                for i, b in enumerate(bullets):
                    para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    para.alignment = PP_ALIGN.LEFT
                    run = para.add_run()
                    run.text = f"• {b}"
                    run.font.size = Pt(20)
                continue

            if layout == "table":
                headers = slide_def.get("headers") or []
                rows = slide_def.get("rows") or []
                if not headers and rows and isinstance(rows[0], dict):
                    headers = list(rows[0].keys())
                ncols = len(headers) if headers else (
                    len(rows[0]) if rows and isinstance(rows[0], list) else 1
                )
                nrows = len(rows) + (1 if headers else 0)
                if nrows < 1:
                    nrows = 1
                table_shape = slide.shapes.add_table(
                    nrows,
                    max(1, ncols),
                    Inches(0.5),
                    Inches(1.5),
                    Inches(12.3),
                    Inches(5.5),
                )
                table = table_shape.table
                if headers:
                    for c, h in enumerate(headers):
                        cell = table.cell(0, c)
                        cell.text = str(h)
                        for run in cell.text_frame.paragraphs[0].runs:
                            run.font.bold = True
                            run.font.size = Pt(14)
                            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(0x30, 0x54, 0x96)
                row_offset = 1 if headers else 0
                for r, row in enumerate(rows):
                    for c in range(ncols):
                        cell = table.cell(r + row_offset, c)
                        if isinstance(row, dict):
                            cell.text = str(row.get(headers[c], "")) if headers else ""
                        elif isinstance(row, list):
                            cell.text = str(row[c]) if c < len(row) else ""
                        else:
                            cell.text = str(row)
                        for run in cell.text_frame.paragraphs[0].runs:
                            run.font.size = Pt(12)
                continue

            if layout == "image":
                b64 = slide_def.get("data_b64") or slide_def.get("image_b64")
                if b64:
                    try:
                        img_bytes = base64.b64decode(b64)
                        img_buf = io.BytesIO(img_bytes)
                        slide.shapes.add_picture(
                            img_buf,
                            Inches(0.5),
                            Inches(1.5),
                            width=Inches(12.3),
                            height=Inches(5.5),
                        )
                        if slide_def.get("caption"):
                            tx = slide.shapes.add_textbox(
                                Inches(0.5),
                                Inches(7.0),
                                Inches(12.3),
                                Inches(0.4),
                            )
                            tf = tx.text_frame
                            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                            run = tf.paragraphs[0].add_run()
                            run.text = str(slide_def["caption"])
                            run.font.size = Pt(12)
                            run.font.italic = True
                    except Exception as err:
                        _logger.warning("PPTX image inválida: %s", err)

        buf = io.BytesIO()
        prs.save(buf)
        slug = self._slugify(fname, "presentacion")
        return {
            "tipo": "pptx",
            **self._store_attachment(
                f"{slug}.pptx",
                buf.getvalue(),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ),
            "nota": "PowerPoint generado. Comparte 'url_descarga' con el usuario.",
        }

    # ==================================================================
    # Matplotlib legacy (fallback)
    # ==================================================================
    def _artifact_chart(self, title, data, chart_kind, fname):
        try:
            import matplotlib

            matplotlib.use("Agg")
            import matplotlib.pyplot as plt
        except ImportError as e:
            raise UserError(_("Instale matplotlib: %s") % e) from e

        fig, ax = plt.subplots(figsize=(6, 4))
        if isinstance(data, list) and data and isinstance(data[0], dict):
            keys = list(data[0].keys())
            if len(keys) >= 2:
                xk, yk = keys[0], keys[1]
                xs = [row.get(xk) for row in data]
                ys = [row.get(yk) for row in data]
                ax.bar(range(len(xs)), ys, tick_label=[str(x) for x in xs])
                ax.set_title(title)
                ax.set_ylabel(yk)
            else:
                ax.text(0.5, 0.5, str(data), ha="center")
        elif isinstance(data, list) and len(data) == 2:
            ax.plot(data[0], data[1])
            ax.set_title(title)
        else:
            ax.text(0.5, 0.5, json.dumps(data, default=str)[:2000], ha="center")
        buf = io.BytesIO()
        fig.savefig(buf, format="png", bbox_inches="tight")
        plt.close(fig)
        slug = self._slugify(fname, "grafico")
        result = self._store_attachment(f"{slug}.png", buf.getvalue(), "image/png")
        base = self.env["ir.config_parameter"].sudo().get_param("web.base.url", "")
        return {
            "tipo": "chart",
            "mensaje_markdown": f"![{title}]({base}/web/image/{result['attachment_id']})",
            **result,
        }
